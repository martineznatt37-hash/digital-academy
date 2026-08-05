# -*- coding: utf-8 -*-
"""Presentación Digital Academy — diseño tech con imágenes."""
import os
import math
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFilter

OUT = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(OUT)
ASSETS = os.path.join(
    os.environ.get('USERPROFILE', ''),
    '.cursor', 'projects', 'c-Users-marti-OneDrive-Desktop-digital-academy', 'assets'
)
SLIDE_W, SLIDE_H = 1920, 1080  # px for backgrounds


def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def find_asset(partial):
    if not os.path.isdir(ASSETS):
        return None
    for f in os.listdir(ASSETS):
        if partial in f:
            return os.path.join(ASSETS, f)
    return None


def img(path):
    return path if path and os.path.isfile(path) else None


def make_tech_bg(name, c1='#030712', c2='#1E3A8A', accent='#2563EB'):
    """Fondo degradado con rejilla y glow — estilo futurista."""
    os.makedirs(os.path.join(OUT, 'slide-assets'), exist_ok=True)
    out = os.path.join(OUT, 'slide-assets', f'{name}.png')
    if os.path.isfile(out):
        return out

    im = Image.new('RGB', (SLIDE_W, SLIDE_H), c1)
    draw = ImageDraw.Draw(im)
    r1, g1, b1 = int(c1[1:3], 16), int(c1[3:5], 16), int(c1[5:7], 16)
    r2, g2, b2 = int(c2[1:3], 16), int(c2[3:5], 16), int(c2[5:7], 16)
    for y in range(SLIDE_H):
        t = y / SLIDE_H
        r = int(r1 + (r2 - r1) * t)
        g = int(g1 + (g2 - g1) * t)
        b = int(b1 + (b2 - b1) * t)
        draw.line([(0, y), (SLIDE_W, y)], fill=(r, g, b))

    ar, ag, ab = int(accent[1:3], 16), int(accent[3:5], 16), int(accent[5:7], 16)
    glow = Image.new('RGBA', (SLIDE_W, SLIDE_H), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    for cx, cy, rad, alpha in [
        (1600, 180, 280, 55), (200, 900, 220, 45), (960, 540, 350, 25)
    ]:
        for i in range(rad, 0, -4):
            a = int(alpha * (1 - i / rad))
            gd.ellipse([cx - i, cy - i, cx + i, cy + i], fill=(ar, ag, ab, a))
    glow = glow.filter(ImageFilter.GaussianBlur(8))
    im = Image.alpha_composite(im.convert('RGBA'), glow).convert('RGB')

    draw = ImageDraw.Draw(im)
    step = 48
    grid_c = (40, 60, 100)
    for x in range(0, SLIDE_W, step):
        draw.line([(x, 0), (x, SLIDE_H)], fill=grid_c, width=1)
    for y in range(0, SLIDE_H, step):
        draw.line([(0, y), (SLIDE_W, y)], fill=grid_c, width=1)

    for x in range(0, SLIDE_W, step * 4):
        for y in range(0, SLIDE_H, step * 4):
            draw.ellipse([x - 2, y - 2, x + 2, y + 2], fill=(ar, ag, ab))

    im.save(out, quality=95)
    return out


def set_bg_image(slide, path):
    slide.shapes.add_picture(path, In(0), In(0), width=In(10), height=In(7.5))


def add_glass_panel(slide, left, top, width, height, alpha=0.88):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb('0F172A')
    sh.fill.transparency = 1 - alpha
    sh.line.color.rgb = rgb('334155')
    sh.line.width = Pt(1)
    return sh


def textbox(slide, left, top, width, height, text, size=18, color='E2E8F0', bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=17, color='CBD5E1', accent='60A5FA'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(8)
        p.level = 0
    return tb


def add_accent_bar(slide, color='2563EB'):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(0), In(0), In(10), In(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(color)
    bar.line.fill.background()


def add_slide_header(slide, title, subtitle=None, accent='2563EB'):
    add_accent_bar(slide, accent)
    textbox(slide, In(0.55), In(0.35), In(8.5), In(0.7), title, size=32, color='FFFFFF', bold=True)
    if subtitle:
        textbox(slide, In(0.55), In(0.95), In(8.5), In(0.45), subtitle, size=14, color='94A3B8')


def add_image_safe(slide, path, left, top, width=None, height=None):
    if not path or not os.path.isfile(path):
        return None
    if width and height:
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(path, left, top, width=width)
    return slide.shapes.add_picture(path, left, top, height=height)


def build_presentation():
    prs = Presentation()
    prs.slide_width = In(10)
    prs.slide_height = In(7.5)
    blank = prs.slide_layouts[6]

    paths = {
        'hero': img(os.path.join(ROOT, 'images', 'site', 'hero.png')) or find_asset('d8e0766'),
        'about': img(os.path.join(ROOT, 'images', 'site', 'about.png')) or find_asset('7744398c'),
        'ai': img(os.path.join(ROOT, 'images', 'site', 'ai-visual.png')),
        'owl': img(os.path.join(ROOT, 'images', 'pets', 'owl.png')),
        'wolf': img(os.path.join(ROOT, 'images', 'pets', 'wolf.png')),
        'dino': img(os.path.join(ROOT, 'images', 'pets', 'dinosaur.png')),
        'ia_course': img(os.path.join(ROOT, 'images', 'courses', 'fundamentos-ia.png')),
        'prog': img(os.path.join(ROOT, 'images', 'courses', 'introduccion-programacion.png')),
        'math_p': img(os.path.join(ROOT, 'images', 'courses', 'matematicas-divertidas.png')),
        'math_prep': img(os.path.join(ROOT, 'images', 'courses', 'matematicas-preparatoria.png')),
    }

    bg_main = make_tech_bg('bg-main', '#020617', '#1E3A8A', '#3B82F6')
    bg_purple = make_tech_bg('bg-purple', '#0F0720', '#4C1D95', '#A855F7')
    bg_cyan = make_tech_bg('bg-cyan', '#042F2E', '#0E7490', '#22D3EE')
    bg_green = make_tech_bg('bg-green', '#052E16', '#14532D', '#22C55E')

    # ── 1 PORTADA ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_main)
    add_glass_panel(s, In(0.45), In(0.55), In(4.8), In(6.4), 0.82)
    textbox(s, In(0.75), In(1.2), In(4.2), In(1.2), 'DIGITAL\nACADEMY', size=44, color='FFFFFF', bold=True)
    textbox(s, In(0.75), In(2.8), In(4.2), In(0.6), 'Aprende. Crece. Transforma tu futuro.', size=20, color='93C5FD')
    textbox(s, In(0.75), In(3.55), In(4.2), In(0.9),
            'Plataforma EdTech con IA · Cursos SEP · Certificados · Exposición 2026', size=13, color='64748B')
    # badges
    for i, (lbl, col) in enumerate([
        ('25 CURSOS', '2563EB'), ('IA 24/7', '7C3AED'), ('100% WEB', '0891B2')
    ]):
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(0.75 + i * 1.45), In(4.8), In(1.25), In(0.42))
        bx.fill.solid()
        bx.fill.fore_color.rgb = rgb(col)
        bx.line.fill.background()
        textbox(s, In(0.75 + i * 1.45), In(4.86), In(1.25), In(0.35), lbl, size=10, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)
    add_image_safe(s, paths['hero'], In(5.35), In(0.65), width=In(4.25))

    # ── 2 PROBLEMA ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, make_tech_bg('bg-red', '#1a0505', '#450a0a', '#EF4444'))
    add_slide_header(s, 'El reto educativo', 'Por qué necesitamos Digital Academy', 'EF4444')
    add_glass_panel(s, In(0.5), In(1.45), In(5.2), In(5.5))
    add_bullets(s, In(0.75), In(1.75), In(4.7), In(4.8), [
        'Contenidos escolares dispersos en internet',
        'Poca alineación con el plan de estudios SEP México',
        'Sin apoyo inmediato fuera del salón de clases',
        'Estudio en línea aburrido → baja motivación',
        'Brecha digital en uso de tecnología',
    ], size=18)
    add_image_safe(s, paths['about'], In(5.85), In(1.55), width=In(3.65))

    # ── 3 SOLUCIÓN ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_main)
    add_slide_header(s, 'Nuestra solución', 'Todo en una sola plataforma web', '2563EB')
    add_glass_panel(s, In(0.5), In(1.45), In(4.6), In(5.5))
    add_bullets(s, In(0.75), In(1.75), In(4.1), In(4.8), [
        'Primaria · Secundaria · Preparatoria',
        '+ Capacitación tecnológica desde cero',
        '25 cursos · 12 lecciones · 3 exámenes',
        'Certificados digitales verificables',
        'Acceso desde cualquier navegador',
    ], size=17)
    # collage cursos
    for i, (p, x, y) in enumerate([
        (paths['math_p'], 5.2, 1.5), (paths['math_prep'], 7.0, 1.5),
        (paths['prog'], 5.2, 3.9), (paths['ia_course'], 7.0, 3.9),
    ]):
        add_image_safe(s, p, In(x), In(y), width=In(1.65))

    # ── 4 FUNCIONES (cards) ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_purple)
    add_slide_header(s, 'Funciones principales', 'Módulos integrados en la plataforma', 'A855F7')
    cards = [
        ('📚', 'Cursos modulares', 'Progreso secuencial\nExámenes por módulo', '2563EB'),
        ('🤖', 'Chat con IA', 'Tutor virtual 24/7\nOpenAI + modo local', '7C3AED'),
        ('👨‍🏫', 'Asesorías', 'Maestros aprobados\nAgenda en línea', '0891B2'),
        ('🐾', 'Mascota virtual', 'Gamificación\nBúho · Lobo · Dino', 'F59E0B'),
        ('🎓', 'Certificados', 'Al completar curso\nID verificable', '22C55E'),
        ('📊', 'Dashboard', 'Gráficas de avance\nLogros y stats', 'EC4899'),
    ]
    for i, (icon, title, desc, col) in enumerate(cards):
        col_i, row_i = i % 3, i // 3
        lx = In(0.45 + col_i * 3.15)
        ty = In(1.55 + row_i * 2.85)
        card = add_glass_panel(s, lx, ty, In(2.95), In(2.55), 0.78)
        textbox(s, lx + In(0.15), ty + In(0.15), In(2.6), In(0.5), icon + '  ' + title, size=16, color='FFFFFF', bold=True)
        textbox(s, lx + In(0.15), ty + In(0.75), In(2.6), In(1.4), desc, size=13, color='94A3B8')
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, ty + In(2.35), In(2.95), In(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = rgb(col)
        strip.line.fill.background()

    # ── 5 FLUJO ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_cyan)
    add_slide_header(s, 'Flujo de aprendizaje', 'Simple de explicar, fácil de usar', '22D3EE')
    steps = [
        ('1', 'Registro', 'Elige tu nivel escolar'),
        ('2', 'Inscripción', 'Selecciona un curso'),
        ('3', 'Estudio', '3 lecciones + examen por módulo'),
        ('4', 'Evaluación', 'Mínimo 8/10 para avanzar'),
        ('5', 'Certificado', 'Al completar los 3 módulos'),
    ]
    for i, (num, title, sub) in enumerate(steps):
        lx = In(0.4 + i * 1.88)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, lx + In(0.55), In(2.0), In(0.75), In(0.75))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb('0891B2')
        circle.line.fill.background()
        textbox(s, lx + In(0.55), In(2.12), In(0.75), In(0.5), num, size=22, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)
        add_glass_panel(s, lx, In(3.0), In(1.75), In(2.2), 0.8)
        textbox(s, lx + In(0.12), In(3.2), In(1.5), In(0.45), title, size=15, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)
        textbox(s, lx + In(0.12), In(3.75), In(1.5), In(1.2), sub, size=12, color='94A3B8', align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, lx + In(1.72), In(2.25), In(0.35), In(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = rgb('22D3EE')
            arr.line.fill.background()
    add_image_safe(s, paths['hero'], In(0.4), In(5.35), width=In(9.2))

    # ── 6 IA ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_purple)
    add_slide_header(s, 'Asistente de Inteligencia Artificial', 'Tu tutor disponible 24/7', '818CF8')
    add_glass_panel(s, In(0.5), In(1.45), In(5.0), In(5.5))
    add_bullets(s, In(0.75), In(1.75), In(4.5), In(4.5), [
        'Widget flotante en toda la plataforma',
        'Responde dudas de matemáticas, ciencias, tecnología…',
        'Contexto de la lección que estás estudiando',
        'OpenAI GPT-4o-mini (opcional)',
        'Modo local sin API como respaldo',
    ], size=17)
    pic = paths['ai'] or paths['ia_course']
    add_image_safe(s, pic, In(5.65), In(1.45), width=In(3.85))

    # ── 7 MASCOTAS ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, make_tech_bg('bg-orange', '#1a0f02', '#78350F', '#F59E0B'))
    add_slide_header(s, 'Gamificación — Mascotas virtuales', 'Aprender también puede ser divertido', 'F59E0B')
    for i, (p, name) in enumerate([(paths['owl'], 'Búho'), (paths['wolf'], 'Lobo'), (paths['dino'], 'Dino')]):
        lx = In(0.55 + i * 3.1)
        add_glass_panel(s, lx, In(1.55), In(2.85), In(5.2), 0.75)
        add_image_safe(s, p, lx + In(0.45), In(1.85), height=In(2.8))
        textbox(s, lx + In(0.2), In(4.85), In(2.45), In(0.5), name, size=22, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)
        textbox(s, lx + In(0.2), In(5.45), In(2.45), In(0.8),
                'Animaciones · Interactiva · Se desbloquea con 10/10', size=11, color='94A3B8', align=PP_ALIGN.CENTER)

    # ── 8 TECNOLOGÍA ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_green)
    add_slide_header(s, 'Stack tecnológico', 'Arquitectura moderna y escalable', '22C55E')
    tech = [
        ('Frontend', 'HTML5 · CSS3 · JavaScript', 'Interfaz responsive'),
        ('Backend', 'Node.js · Express', 'API REST JSON'),
        ('Base de datos', 'SQLite', '19 tablas relacionales'),
        ('Seguridad', 'JWT · bcrypt', 'Roles y sesiones'),
        ('Visualización', 'Chart.js', 'Gráficas de progreso'),
        ('IA', 'OpenAI API', 'Chat contextual'),
    ]
    for i, (layer, stack, note) in enumerate(tech):
        col_i, row_i = i % 3, i // 3
        lx = In(0.45 + col_i * 3.15)
        ty = In(1.55 + row_i * 2.85)
        add_glass_panel(s, lx, ty, In(2.95), In(2.55), 0.82)
        textbox(s, lx + In(0.2), ty + In(0.2), In(2.55), In(0.4), layer, size=13, color='22C55E', bold=True)
        textbox(s, lx + In(0.2), ty + In(0.65), In(2.55), In(0.55), stack, size=16, color='FFFFFF', bold=True)
        textbox(s, lx + In(0.2), ty + In(1.35), In(2.55), In(0.8), note, size=12, color='94A3B8')

    # ── 9 ROLES ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_main)
    add_slide_header(s, 'Roles del sistema', 'Tres tipos de usuario', '2563EB')
    roles = [
        ('👤 Estudiante', ['Cursos y exámenes', 'Chat IA y mascota', 'Certificados', 'Asesorías'], '2563EB'),
        ('👨‍🏫 Maestro', ['Registro + aprobación', 'Panel de sesiones', 'Datos del alumno', 'Notificaciones'], '7C3AED'),
        ('🔐 Administrador', ['Aprobar maestros', 'Confirmar asesorías', 'Estadísticas', 'Control total'], 'F59E0B'),
    ]
    for i, (title, items, col) in enumerate(roles):
        lx = In(0.45 + i * 3.15)
        add_glass_panel(s, lx, In(1.55), In(2.95), In(5.2), 0.8)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, In(1.55), In(2.95), In(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(col)
        bar.line.fill.background()
        textbox(s, lx + In(0.15), In(1.85), In(2.65), In(0.5), title, size=18, color='FFFFFF', bold=True)
        add_bullets(s, lx + In(0.15), In(2.55), In(2.65), In(3.8), ['• ' + x for x in items], size=14)

    # ── 10 NÚMEROS ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_cyan)
    add_slide_header(s, 'Resultados en números', 'Impacto del proyecto', '06B6D4')
    stats = [
        ('25', 'Cursos activos'), ('300', 'Lecciones totales'), ('30+', 'Endpoints API'),
        ('19', 'Tablas en BD'), ('3', 'Roles de usuario'), ('100%', 'Funcional'),
    ]
    for i, (num, lbl) in enumerate(stats):
        col_i, row_i = i % 3, i // 3
        lx = In(0.45 + col_i * 3.15)
        ty = In(1.65 + row_i * 2.75)
        add_glass_panel(s, lx, ty, In(2.95), In(2.35), 0.85)
        textbox(s, lx + In(0.1), ty + In(0.35), In(2.75), In(0.9), num, size=48, color='22D3EE', bold=True, align=PP_ALIGN.CENTER)
        textbox(s, lx + In(0.1), ty + In(1.45), In(2.75), In(0.5), lbl, size=14, color='CBD5E1', align=PP_ALIGN.CENTER)

    # ── 11 DEMO ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, bg_main)
    add_slide_header(s, 'Demo en vivo', 'Prueba la plataforma ahora', '2563EB')
    add_glass_panel(s, In(0.5), In(1.45), In(5.5), In(5.5))
    add_bullets(s, In(0.75), In(1.85), In(5.0), In(4.5), [
        '🌐  http://localhost:3001',
        '',
        '👤 Estudiante: demo@digitalacademy.com / demo1234',
        '🔐 Admin: admin@gmail.com / 1234',
        '',
        'Mostrar: Catálogo → Curso → Examen → Chat IA → Mascota',
    ], size=17)
    add_image_safe(s, paths['hero'], In(6.15), In(1.55), width=In(3.35))

    # ── 12 CIERRE ──
    s = prs.slides.add_slide(blank)
    set_bg_image(s, make_tech_bg('bg-close', '#1E3A8A', '#020617', '#60A5FA'))
    add_glass_panel(s, In(1.2), In(1.8), In(7.6), In(3.9), 0.72)
    textbox(s, In(1.5), In(2.3), In(7.0), In(1.0), '¡Gracias!', size=52, color='FFFFFF', bold=True, align=PP_ALIGN.CENTER)
    textbox(s, In(1.5), In(3.4), In(7.0), In(0.6),
            'Digital Academy — Educación + Tecnología + Inteligencia Artificial',
            size=20, color='93C5FD', align=PP_ALIGN.CENTER)
    textbox(s, In(1.5), In(4.2), In(7.0), In(0.5),
            '¿Preguntas?', size=16, color='64748B', align=PP_ALIGN.CENTER)
    # mini pets footer
    for i, p in enumerate([paths['owl'], paths['wolf'], paths['dino']]):
        add_image_safe(s, p, In(3.2 + i * 1.2), In(5.5), height=In(1.1))

    path = os.path.join(OUT, 'Presentacion_Digital_Academy.pptx')
    prs.save(path)
    print(f'Generado: {path} ({len(prs.slides)} diapositivas)')
    return path


if __name__ == '__main__':
    build_presentation()
