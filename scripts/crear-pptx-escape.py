# -*- coding: utf-8 -*-
"""Genera diapositivas de exposición Escape del Hacker (estilo ciber/neon)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import copy

OUT = r"C:\Users\marti\OneDrive\Desktop\digital-academy\docs\Presentacion-Escape-Del-Hacker.pptx"
OUT2 = r"C:\Users\marti\OneDrive\Desktop\Escape-Del-Hacker\docs\Presentacion-Escape-Del-Hacker.pptx"

CYAN = RGBColor(0x00, 0xFF, 0xFF)
PINK = RGBColor(0xFF, 0x44, 0xAA)
YELLOW = RGBColor(0xFF, 0xDD, 0x66)
WHITE = RGBColor(0xEE, 0xEE, 0xEE)
GRAY = RGBColor(0x99, 0xAA, 0xBB)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x05, 0x0A, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_run(run, size=20, bold=False, color=WHITE, font="Consolas"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def paint_bg(slide, rgb=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_bar(slide, top=False):
    # franja neón superior/inferior
    y = 0 if top else Inches(7.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()


def add_title(slide, text, y=0.35, size=36, color=CYAN):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=True, color=color)
    return box


def add_body(slide, lines, y=1.5, size=20, color=WHITE):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, color=color)
    return box


def new_slide(title, lines, subtitle=None):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    paint_bg(slide)
    add_bar(slide, top=True)
    add_bar(slide, top=False)
    add_title(slide, title)
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = subtitle
        set_run(run, size=16, color=YELLOW)
        add_body(slide, lines, y=1.7)
    else:
        add_body(slide, lines, y=1.5)
    # marca de agua
    mark = slide.shapes.add_textbox(Inches(10.2), Inches(6.9), Inches(2.8), Inches(0.35))
    r = mark.text_frame.paragraphs[0].add_run()
    r.text = "ESCAPE DEL HACKER"
    set_run(r, size=10, color=GRAY)
    return slide


# ---- Diapositivas ----

# Portada
s = prs.slides.add_slide(prs.slide_layouts[6])
paint_bg(s)
add_bar(s, True)
add_bar(s, False)
# panel
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.6), Inches(10.3), Inches(4.2))
panel.fill.solid()
panel.fill.fore_color.rgb = RGBColor(0x00, 0x18, 0x22)
panel.line.color.rgb = CYAN
panel.line.width = Pt(2)
t = s.shapes.add_textbox(Inches(1.8), Inches(2.0), Inches(9.7), Inches(3.5))
tf = t.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "ESCAPE DEL HACKER 2.0"
set_run(r, size=40, bold=True, color=CYAN)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Juego web de infiltración y ciberseguridad"
set_run(r2, size=20, color=WHITE)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "Universidad Azteca · Ingeniería en Sistemas Computacionales"
set_run(r3, size=14, color=YELLOW)
p4 = tf.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run()
r4.text = "Interfaz Gráfica · Demostración del proyecto"
set_run(r4, size=14, color=GRAY)

new_slide(
    ">> AGENDA",
    [
        "01  Problema y objetivo",
        "02  ¿Qué es Escape del Hacker?",
        "03  Interfaz y experiencia de usuario",
        "04  Funcionalidades del sistema",
        "05  Arquitectura y escalabilidad",
        "06  Tecnologías y despliegue",
        "07  Demostración en vivo",
        "08  Beneficios y conclusiones",
    ],
)

new_slide(
    ">> PROBLEMA Y OBJETIVO",
    [
        "Problema: conceptos de malware suelen verse solo en teoría.",
        "Objetivo: crear una interfaz jugable que ensene ciberseguridad",
        "de forma visual, atractiva y accesible en cualquier dispositivo.",
        "",
        "Meta de la rúbrica:",
        "• UI/UX atractiva y multiplataforma",
        "• Funciones clave operando (cuenta, juego, progreso)",
        "• Software escalable y documentado",
    ],
)

new_slide(
    ">> ¿QUÉ ES EL JUEGO?",
    [
        "El jugador es teletransportado a un sistema hackeado.",
        "Debe escapar antes de que el hacker lo encuentre.",
        "",
        "• 10 niveles · 1 malware por nivel",
        "• Recoger llaves · evitar enemigos · llegar a la puerta EXIT",
        "• Personajes con habilidades (escudo, sigilo, ralentizar...)",
        "• Jefes en niveles avanzados",
        "• Cuenta con usuario y contraseña para continuar en otro dispositivo",
    ],
    subtitle="Concepto: infiltración + aprendizaje lúdico",
)

new_slide(
    ">> INTERFAZ GRÁFICA (UI/UX)",
    [
        "Diseño oscuro con neón cian/rosa (estética hacker).",
        "Navegación fluida: Login → Menú → Personaje → Intro → Juego.",
        "",
        "PC: teclado (flechas/WASD + E habilidad + Espacio ataque).",
        "Móvil/tablet: D-pad táctil + botón de habilidad.",
        "",
        "Accesible en cualquier dispositivo mediante enlace Netlify.",
        "PWA: se puede agregar a pantalla de inicio como app.",
    ],
)

new_slide(
    ">> FUNCIONALIDADES CLAVE",
    [
        "1. Registro e inicio de sesión",
        "2. Selección de personaje y habilidades",
        "3. Gameplay completo (niveles, llaves, enemigos, jefes)",
        "4. Guardado de récord y nivel máximo",
        "5. Sincronización por API (misma cuenta en varios dispositivos)",
        "6. Publicación web (Netlify) + servidor (Render)",
        "",
        "Equivalente en la rúbrica a “funciones principales operando correctamente”.",
    ],
)

new_slide(
    ">> ARQUITECTURA ESCALABLE",
    [
        "Cliente (HTML/JS/Canvas)  →  API Express (/api/escape)  →  SQLite",
        "",
        "Módulos separados:",
        "• Frontend del juego (escape-public)",
        "• Backend de cuentas (server/routes/escape.js)",
        "• Datos configurables: niveles[], personajes[], temas{}",
        "",
        "Para crecer: agregar nivel/personaje/ruta API sin rehacer el núcleo.",
        "Ejemplos futuros: ranking global, tienda de skins, más jefes.",
    ],
)

new_slide(
    ">> TECNOLOGÍAS",
    [
        "Frontend: HTML5 · CSS3 · JavaScript · Canvas 2D",
        "Backend: Node.js · Express · SQLite · bcrypt · JWT",
        "Deploy: GitHub · Netlify (juego) · Render (API)",
        "",
        "Seguridad básica de cuentas:",
        "• Contraseñas hasheadas (bcrypt)",
        "• Sesión con token JWT",
    ],
)

new_slide(
    ">> DESPLIEGUE",
    [
        "1. Código en GitHub (Escape-Del-Hacker)",
        "2. Frontend en Netlify → enlace público",
        "3. API en Render → https://escape-del-hacker.onrender.com",
        "4. Conexión:",
        "   sitio.netlify.app/?api=https://escape-del-hacker.onrender.com/api/escape",
        "",
        "Resultado: jugar desde PC, celular o tablet con la misma cuenta.",
    ],
)

new_slide(
    ">> DEMOSTRACIÓN EN VIVO",
    [
        "1. Abrir el link de Netlify (con API)",
        "2. Registrar / Entrar",
        "3. Elegir personaje",
        "4. Completar o avanzar un nivel",
        "5. Mostrar controles en celular",
        "6. (Opcional) Entrar con la misma cuenta en otro dispositivo",
        "",
        "Enfoque: beneficios reales para el usuario final.",
    ],
)

new_slide(
    ">> BENEFICIOS PARA EL USUARIO",
    [
        "• Aprende malware jugando, no solo leyendo",
        "• No necesita instalar un motor pesado: abre el navegador",
        "• Continúa su progreso en otro aparato",
        "• Experiencia visual atractiva y motivadora",
        "• Disponible en la nube (compartible por link)",
    ],
)

new_slide(
    ">> CONCLUSIÓN",
    [
        "Se entregó una interfaz gráfica completa, jugable y publicada.",
        "Cumple UI/UX multiplataforma, funcionalidad, escalabilidad",
        "y documentación técnica para la exposición.",
        "",
        "ESCAPE DEL HACKER 2.0",
        "Infiltración · Aprendizaje · Escapatoria",
        "",
        "¿Preguntas?",
    ],
)

prs.save(OUT)
try:
    import os
    os.makedirs(os.path.dirname(OUT2), exist_ok=True)
    prs.save(OUT2)
except Exception as e:
    print("OUT2 skip", e)
print("OK", OUT)
